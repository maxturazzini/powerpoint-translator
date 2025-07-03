from typing import Dict, List, Tuple, Optional, Any
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from pptx.text.text import _Run
from pptx.dml.color import ColorFormat
from formatting import TextRunFormatting

class VisualComparator:
    """Compare visual formatting between original and translated presentations"""
    
    @staticmethod
    def _get_color_info(font) -> Optional[Dict[str, Any]]:
        """Safely extract color information from a font"""
        if not font.color:
            return None
            
        color_info = {
            'type': font.color.type
        }
        
        # Handle different color types
        if hasattr(font.color, 'rgb') and font.color.rgb is not None:
            color_info['rgb'] = font.color.rgb
        if hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
            color_info['theme_color'] = font.color.theme_color
        if hasattr(font.color, 'brightness') and font.color.brightness is not None:
            color_info['brightness'] = font.color.brightness
            
        return color_info if any(k != 'type' for k in color_info) else None
    
    @staticmethod
    def compare_runs(
        original_run: _Run,
        translated_run: _Run
    ) -> List[str]:
        """Compare formatting between two text runs"""
        differences = []
        
        # Compare font properties
        if original_run.font.name != translated_run.font.name:
            differences.append(
                f"Font name mismatch: {original_run.font.name} vs "
                f"{translated_run.font.name}"
            )
            
        if original_run.font.size != translated_run.font.size:
            differences.append(
                f"Font size mismatch: {original_run.font.size} vs "
                f"{translated_run.font.size}"
            )
            
        # Compare text formatting
        for attr in ['bold', 'italic', 'underline']:
            orig_value = getattr(original_run.font, attr)
            trans_value = getattr(translated_run.font, attr)
            if orig_value != trans_value:
                differences.append(
                    f"{attr.capitalize()} formatting mismatch: "
                    f"{orig_value} vs {trans_value}"
                )
                
        # Compare colors safely
        orig_color = VisualComparator._get_color_info(original_run.font)
        trans_color = VisualComparator._get_color_info(translated_run.font)
        
        if orig_color != trans_color:
            differences.append(
                f"Color mismatch: {orig_color} vs {trans_color}"
            )
            
        return differences

    @staticmethod
    def generate_formatting_report(
        shape: BaseShape,
        is_original: bool = True
    ) -> Dict[str, List[TextRunFormatting]]:
        """Generate a detailed formatting report for a shape"""
        report = {}
        prefix = "Original" if is_original else "Translated"
        
        if hasattr(shape, 'text_frame'):
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                # Create a list to store formats for this paragraph
                para_formats = []
                
                # Process each run in the paragraph
                for run in paragraph.runs:
                    # Get color information safely
                    color_info = VisualComparator._get_color_info(run.font)
                    
                    # Get language ID safely
                    try:
                        language_id = run.font.language_id
                    except (ValueError, AttributeError):
                        language_id = None
                    
                    # Get spacing safely
                    try:
                        spacing = run.font.spacing
                    except AttributeError:
                        spacing = None
                    
                    # Get paragraph properties
                    alignment = getattr(paragraph, 'alignment', None)
                    level = getattr(paragraph, 'level', None)
                    
                    # Get line spacing and paragraph spacing
                    line_spacing = None
                    space_before = None
                    space_after = None
                    if hasattr(paragraph._element, 'pPr'):
                        pPr = paragraph._element.pPr
                        if hasattr(pPr, 'lnSpc'):
                            line_spacing = pPr.lnSpc
                        if hasattr(pPr, 'spcBef'):
                            space_before = pPr.spcBef
                        if hasattr(pPr, 'spcAft'):
                            space_after = pPr.spcAft
                    
                    formatting = TextRunFormatting(
                        font_name=run.font.name,
                        font_size=run.font.size,
                        bold=run.font.bold,
                        italic=run.font.italic,
                        underline=run.font.underline,
                        color_format=color_info,
                        language_id=language_id,
                        spacing=spacing,
                        alignment=alignment,
                        line_spacing=line_spacing,
                        space_before=space_before,
                        space_after=space_after,
                        level=level
                    )
                    para_formats.append(formatting)
                
                # Only add to report if there are formats
                if para_formats:
                    report[f"{prefix} Paragraph {para_idx + 1}"] = para_formats
                
        return report

    @staticmethod
    def compare_shapes(
        original_shape: BaseShape,
        translated_shape: BaseShape
    ) -> Tuple[List[str], Dict[str, List[TextRunFormatting]]]:
        """
        Compare formatting between two shapes.
        Returns (differences, formatting_report)
        """
        differences = []
        
        # Generate formatting reports
        try:
            orig_report = VisualComparator.generate_formatting_report(
                original_shape,
                is_original=True
            )
            trans_report = VisualComparator.generate_formatting_report(
                translated_shape,
                is_original=False
            )
        except Exception as e:
            differences.append(f"Error generating formatting report: {str(e)}")
            return differences, {}
        
        # Compare text frame properties
        if hasattr(original_shape, 'text_frame') and hasattr(translated_shape, 'text_frame'):
            orig_tf = original_shape.text_frame
            trans_tf = translated_shape.text_frame
            
            # Compare paragraph count
            if len(orig_tf.paragraphs) != len(trans_tf.paragraphs):
                differences.append(
                    f"Paragraph count mismatch: {len(orig_tf.paragraphs)} vs "
                    f"{len(trans_tf.paragraphs)}"
                )
                
            # Compare each paragraph
            for para_idx, (orig_para, trans_para) in enumerate(
                zip(orig_tf.paragraphs, trans_tf.paragraphs)
            ):
                # Compare run count
                if len(orig_para.runs) != len(trans_para.runs):
                    differences.append(
                        f"Paragraph {para_idx + 1}: Run count mismatch: "
                        f"{len(orig_para.runs)} vs {len(trans_para.runs)}"
                    )
                    
                # Compare each run
                for run_idx, (orig_run, trans_run) in enumerate(
                    zip(orig_para.runs, trans_para.runs)
                ):
                    try:
                        run_diffs = VisualComparator.compare_runs(orig_run, trans_run)
                        if run_diffs:
                            differences.extend([
                                f"Paragraph {para_idx + 1}, Run {run_idx + 1}: {diff}"
                                for diff in run_diffs
                            ])
                    except Exception as e:
                        differences.append(
                            f"Error comparing runs in paragraph {para_idx + 1}, "
                            f"run {run_idx + 1}: {str(e)}"
                        )
                        
        # Combine reports
        formatting_report = {**orig_report, **trans_report}
        
        return differences, formatting_report

    @staticmethod
    def generate_comparison_summary(
        differences: List[str],
        formatting_report: Dict[str, List[TextRunFormatting]]
    ) -> str:
        """Generate a human-readable comparison summary"""
        summary = []
        
        if differences:
            summary.append("Formatting Differences Found:")
            for diff in differences:
                summary.append(f"- {diff}")
        else:
            summary.append("No formatting differences found.")
            
        summary.append("\nFormatting Details:")
        for location, formats in formatting_report.items():
            summary.append(f"\n{location}:")
            for idx, fmt in enumerate(formats, 1):
                summary.append(f"  Format {idx}:")
                for attr, value in fmt.to_dict().items():
                    if value is not None:
                        summary.append(f"    {attr}: {value}")
                        
        return "\n".join(summary)