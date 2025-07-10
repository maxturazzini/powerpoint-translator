from typing import List, Dict, Any, Optional
from pptx.presentation import Presentation
from pptx.shapes.base import BaseShape
from formatting import FormattingManager
from processors import TextProcessor

class FormatValidator:
    """Validate formatting preservation between original and translated presentations"""
    
    def __init__(self, formatting_manager: FormattingManager):
        self.formatting_manager = formatting_manager
        self.text_processor = TextProcessor()
        self.warnings: List[str] = []
        
    def validate_presentation(
        self,
        original_prs: Presentation,
        translated_prs: Presentation
    ) -> List[str]:
        """
        Validate formatting preservation between two presentations.
        Returns a list of warning messages.
        """
        self.warnings = []
        
        if len(original_prs.slides) != len(translated_prs.slides):
            self.warnings.append(
                f"Slide count mismatch: Original={len(original_prs.slides)}, "
                f"Translated={len(translated_prs.slides)}"
            )
            
        # Validate each slide
        for slide_idx, (orig_slide, trans_slide) in enumerate(
            zip(original_prs.slides, translated_prs.slides)
        ):
            self._validate_slide(slide_idx + 1, orig_slide, trans_slide)
            
        return self.warnings
        
    def _validate_slide(self, slide_num: int, orig_slide, trans_slide) -> None:
        """Validate formatting preservation for a single slide"""
        if len(orig_slide.shapes) != len(trans_slide.shapes):
            self.warnings.append(
                f"Slide {slide_num}: Shape count mismatch: "
                f"Original={len(orig_slide.shapes)}, "
                f"Translated={len(trans_slide.shapes)}"
            )
            
        # Validate each shape
        for shape_idx, (orig_shape, trans_shape) in enumerate(
            zip(orig_slide.shapes, trans_slide.shapes)
        ):
            self._validate_shape(
                f"Slide {slide_num}, Shape {shape_idx + 1}",
                orig_shape,
                trans_shape
            )
            
        # Validate notes
        if orig_slide.has_notes_slide != trans_slide.has_notes_slide:
            self.warnings.append(
                f"Slide {slide_num}: Notes presence mismatch"
            )
        elif orig_slide.has_notes_slide:
            self._validate_notes(
                slide_num,
                orig_slide.notes_slide,
                trans_slide.notes_slide
            )
            
    def _validate_shape(
        self,
        location: str,
        orig_shape: BaseShape,
        trans_shape: BaseShape
    ) -> None:
        """Validate formatting preservation for a single shape"""
        if orig_shape.shape_type != trans_shape.shape_type:
            self.warnings.append(
                f"{location}: Shape type mismatch: "
                f"Original={orig_shape.shape_type}, "
                f"Translated={trans_shape.shape_type}"
            )
            return
            
        # Validate text content and formatting
        orig_content = self.text_processor.extract_shape_content(orig_shape)
        trans_content = self.text_processor.extract_shape_content(trans_shape)
        
        content_warnings = self.text_processor.verify_translation_integrity(
            orig_content,
            trans_content
        )
        
        for warning in content_warnings:
            self.warnings.append(f"{location}: {warning}")
            
        # Validate formatting preservation
        format_warnings = self.formatting_manager.validate_formatting(
            str(orig_shape.shape_id),
            str(trans_shape.shape_id)
        )
        
        for warning in format_warnings:
            self.warnings.append(f"{location}: {warning}")
            
    def _validate_notes(self, slide_num: int, orig_notes, trans_notes) -> None:
        """Validate formatting preservation for slide notes"""
        orig_content = self.text_processor.extract_notes_content(orig_notes)
        trans_content = self.text_processor.extract_notes_content(trans_notes)
        
        if bool(orig_content) != bool(trans_content):
            self.warnings.append(
                f"Slide {slide_num}: Notes content presence mismatch"
            )
            return
            
        if orig_content:
            # Validate notes formatting
            format_warnings = self.formatting_manager.validate_formatting(
                "notes",
                "notes"
            )
            
            for warning in format_warnings:
                self.warnings.append(f"Slide {slide_num} Notes: {warning}")
                
    def get_validation_summary(self) -> Dict[str, Any]:
        """
        Generate a summary of validation results.
        Returns a dictionary with validation statistics.
        """
        return {
            'total_warnings': len(self.warnings),
            'format_warnings': len([w for w in self.warnings if 'format' in w.lower()]),
            'content_warnings': len([w for w in self.warnings if 'content' in w.lower()]),
            'structure_warnings': len([w for w in self.warnings if 'mismatch' in w.lower()]),
            'warnings': self.warnings
        }
