import sys
import os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext
import threading
import webbrowser
import platform
import subprocess
from translate_powerpoint import PowerPointTranslator, INSTRUCTIONS
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from processors.text_processor import TextProcessor


def export_presentation_to_markdown(pptx_path: str):
    """Export all slides of a presentation to a Markdown file."""
    prs = Presentation(pptx_path)
    md_lines = []
    summary = []
    for idx, slide in enumerate(prs.slides, 1):
        md_lines.append(f"# Slide {idx}")

        title = ""
        if slide.shapes.title:
            title = TextProcessor.get_text_frame_content(slide.shapes.title.text_frame)
            md_lines.append(f"## {title}")
        else:
            md_lines.append("## ")

        main_parts = []
        all_texts = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = TextProcessor.get_text_frame_content(shape.text_frame)
            if not text.strip():
                continue
            if shape.is_placeholder and getattr(shape.placeholder_format, "type", None) in (
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.CENTER_BODY,
            ):
                main_parts.append(text)
            all_texts.append(text)

        main_text = " ".join(main_parts)
        md_lines.append(f"### {main_text}")
        md_lines.append(f"#### {' '.join(all_texts)}")

        if slide.has_notes_slide:
            notes = TextProcessor.get_text_frame_content(slide.notes_slide.notes_text_frame)
            md_lines.append(f"##### {notes}")

        md_lines.append("")
        summary.append(f"{os.path.basename(pptx_path)} | {idx} | {title} | {main_text}")

    md_path = os.path.splitext(pptx_path)[0] + ".md"
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))

    return summary

# Load OpenAI API key from .env using python-dotenv
# Requires: pip install python-dotenv
try:
    from dotenv import load_dotenv
    # Look for .env in current directory first, then parent directory
    current_dir = os.path.dirname(os.path.abspath(__file__))
    env_paths = [
        os.path.join(current_dir, '.env'),  # Current directory
        os.path.join(os.path.dirname(current_dir), '.env')  # Parent directory
    ]
    
    for env_path in env_paths:
        if os.path.exists(env_path):
            load_dotenv(env_path)
            break
    
    # Try different possible environment variable names
    ENV_API_KEY = os.getenv('openai_api_key') or os.getenv('OPENAI_API_KEY')
except ImportError:
    ENV_API_KEY = None

class TranslatorGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Translator GUI")
        self.api_key = ENV_API_KEY or ""
        self.model = "gpt-4o-mini"
        self.translated_path = None

        # File selection
        self.input_path_var = tk.StringVar()
        self.output_path_var = tk.StringVar()
        self.prompt_var = tk.StringVar(value=INSTRUCTIONS.strip())

        tk.Label(root, text="Select PowerPoint file to translate:").pack(anchor="w", padx=10, pady=(10,0))
        file_frame = tk.Frame(root)
        file_frame.pack(fill="x", padx=10)
        tk.Entry(file_frame, textvariable=self.input_path_var, width=60).pack(side="left", expand=True, fill="x")
        tk.Button(file_frame, text="Browse", command=self.browse_file).pack(side="left", padx=5)

        tk.Label(root, text="Destination file (auto-generated):").pack(anchor="w", padx=10, pady=(10,0))
        out_frame = tk.Frame(root)
        out_frame.pack(fill="x", padx=10)
        out_entry = tk.Entry(out_frame, textvariable=self.output_path_var, width=60, state="readonly")
        out_entry.pack(side="left", expand=True, fill="x")

        tk.Label(root, text="System Prompt (edit if needed):").pack(anchor="w", padx=10, pady=(10,0))
        self.prompt_text = scrolledtext.ScrolledText(root, height=12, wrap=tk.WORD)
        self.prompt_text.pack(fill="both", expand=True, padx=10)
        self.prompt_text.insert(tk.END, INSTRUCTIONS.strip())

        # Translate button
        self.translate_btn = tk.Button(root, text="Translate", command=self.start_translation)
        self.translate_btn.pack(pady=10)

        # Open file button (hidden by default)
        self.open_btn = tk.Button(root, text="Open Translated File", command=self.open_translated_file)
        self.open_btn.pack(pady=(0,10))
        self.open_btn.pack_forget()

        # Export to Markdown
        self.export_btn = tk.Button(root, text="Export Slides to Markdown", command=self.export_slides)
        self.export_btn.pack(pady=(0,10))

        # Status
        self.status_var = tk.StringVar()
        
        # Set initial status based on API key availability
        if self.api_key and self.api_key.strip():
            self.status_var.set("✅ API key loaded - Ready to translate")
        else:
            self.status_var.set("⚠️ No API key found - Will prompt during translation")
            
        tk.Label(root, textvariable=self.status_var, fg="blue").pack(pady=(0,10))

    def browse_file(self):
        path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
        if path:
            self.input_path_var.set(path)
            base, ext = os.path.splitext(path)
            dest = f"{base}_Translated{ext}"
            self.output_path_var.set(dest)
            self.translated_path = dest
            self.open_btn.pack_forget()

    def start_translation(self):
        input_path = self.input_path_var.get()
        output_path = self.output_path_var.get()
        prompt = self.prompt_text.get("1.0", tk.END).strip()
        if not input_path or not output_path:
            messagebox.showerror("Error", "Please select a PowerPoint file.")
            return
        if not self.api_key or not self.api_key.strip():
            self.api_key = self.ask_api_key()
            if not self.api_key or not self.api_key.strip():
                return
        self.status_var.set("Translating... Please wait.")
        self.translate_btn.config(state=tk.DISABLED)
        self.open_btn.pack_forget()
        threading.Thread(target=self.run_translation, args=(input_path, output_path, prompt), daemon=True).start()

    def ask_api_key(self):
        # Simple dialog for API key
        return simple_input_dialog(self.root, "Enter OpenAI API Key:")

    def run_translation(self, input_path, output_path, prompt):
        try:
            translator = PowerPointTranslator(
                api_key=self.api_key,
                model=self.model,
                translate_notes=True,
                skip_hidden_slides=True
            )
            # Patch the INSTRUCTIONS for this run
            global INSTRUCTIONS
            INSTRUCTIONS = prompt
            translator.translate_presentation(input_path, output_path)
            self.status_var.set(f"Done! Saved to: {output_path}")
            self.translated_path = output_path
            self.open_btn.pack(pady=(0,10))
            messagebox.showinfo("Success", f"Translation complete!\nSaved to: {output_path}")
        except Exception as e:
            self.status_var.set(f"Error: {e}")
            messagebox.showerror("Error", str(e))
        finally:
            self.translate_btn.config(state=tk.NORMAL)

    def open_translated_file(self):
        if self.translated_path and os.path.exists(self.translated_path):
            system = platform.system()
            path = os.path.abspath(self.translated_path)
            if system == "Darwin":
                subprocess.run(["open", "-R", path])
            elif system == "Windows":
                os.startfile(os.path.normpath(path))
            else:
                subprocess.run(["xdg-open", path])
        else:
            messagebox.showerror("Error", "Translated file not found.")

    def export_slides(self):
        """Export slides to Markdown from a file or folder."""
        file_path = filedialog.askopenfilename(title="Select PowerPoint file", filetypes=[("PowerPoint files", "*.pptx")])
        paths = []
        if file_path:
            paths = [file_path]
        else:
            folder = filedialog.askdirectory(title="Select folder containing PPTX files")
            if not folder:
                return
            for name in os.listdir(folder):
                if name.lower().endswith(".pptx"):
                    paths.append(os.path.join(folder, name))

        if not paths:
            messagebox.showinfo("No files", "No PowerPoint files selected.")
            return

        summary_lines = []
        for p in paths:
            self.status_var.set(f"Exporting {os.path.basename(p)}")
            self.root.update_idletasks()
            try:
                summary_lines.extend(export_presentation_to_markdown(p))
            except Exception as e:
                messagebox.showerror("Error", f"Failed to export {p}: {e}")

        if summary_lines:
            with open("AllTheSlides.md", "w", encoding="utf-8") as f:
                for line in summary_lines:
                    f.write(line + "\n")
        messagebox.showinfo("Done", "Markdown export completed")
        self.status_var.set("Markdown export completed")

def simple_input_dialog(parent, prompt):
    dialog = tk.Toplevel(parent)
    dialog.title("Input Required")
    tk.Label(dialog, text=prompt).pack(padx=10, pady=10)
    entry = tk.Entry(dialog, width=50, show="*")
    entry.pack(padx=10, pady=5)
    entry.focus_set()
    result = []
    def on_ok():
        result.append(entry.get())
        dialog.destroy()
    tk.Button(dialog, text="OK", command=on_ok).pack(pady=10)
    dialog.transient(parent)
    dialog.grab_set()
    parent.wait_window(dialog)
    return result[0] if result else None

if __name__ == "__main__":
    root = tk.Tk()
    app = TranslatorGUI(root)
    root.mainloop()
