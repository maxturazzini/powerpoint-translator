"""
PowerPoint Translator V5
Implements comprehensive text formatting preservation during translation.

Features:
- Complete formatting preservation
- Enhanced text processing
- Comprehensive validation
- Visual comparison tools
"""

from openai import OpenAI
import logging
import sys
from typing import Optional, Callable
from pptx import Presentation
from formatting import FormattingManager
from processors import TextProcessor
from processors.enhanced_shape_processor import EnhancedShapeProcessor
from validation import FormatValidator, VisualComparator

# Configure logging to both file and console
logger = logging.getLogger('ppt_translator')
logger.setLevel(logging.INFO)

# Clear any existing handlers
logger.handlers = []

# File handler
file_handler = logging.FileHandler('ppt_translator.log')
file_handler.setLevel(logging.INFO)
file_handler.setFormatter(
    logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
)
logger.addHandler(file_handler)

# Console handler
console_handler = logging.StreamHandler(sys.stdout)
console_handler.setLevel(logging.INFO)
console_handler.setFormatter(
    logging.Formatter('%(message)s')  # Simplified console format
)
logger.addHandler(console_handler)

# Translation instructions
INSTRUCTIONS = """
You are a professional translator specialising in business- and technical-oriented PowerPoint presentations.
Your assignment is to translate all text written in English (SOURCE_LANG) into Italian (TARGET_LANG).

### Task
* Translate every segment of text that is in SOURCE_LANG to TARGET_LANG taking into account the context of the text (i.e. technical terms, names, numbers, etc.).

### Translation guidelines

1. Produce clear, fluent, professional TARGET_LANG.
2. Preserve meaning, tone, and intent.
3. Keep proper names, acronyms, technical terms, and numbers exactly as written.
4. Retain **all** original formatting — line breaks, spacing, punctuation, and special characters.
5. Provide **only** the translated text; add **no explanations or comments**.
6. Use consistent terminology throughout.

### Example
**Input**
Welcome to the
final presentation.
**Output**
Benvenuti alla
presentazione finale.
"""

exINSTRUCTIONS = """
You are a professional Italian to English translator. Your task is to translate PowerPoint content from Italian to English.

IMPORTANT: You must ALWAYS translate any Italian text to English. Never leave Italian text untranslated.

Guidelines:
1. ALWAYS translate Italian text to clear, professional business English
2. Maintain the original meaning and context
3. Keep technical terms, names, and numbers unchanged
4. If text is already in English, return it unchanged
5. Preserve all formatting, spaces, and special characters
6. Provide only the translated text, no explanations
7. Ensure natural, fluent English phrasing
8. Use consistent terminology throughout
9. Preserve line breaks and spacing exactly as in the original text

Example:
Input (Italian): "Benvenuti alla\npresentazione"
Output (English): "Welcome to the\npresentation"

Remember: ALL Italian text MUST be translated to English.
"""

class PowerPointTranslator:
    """PowerPoint translator with comprehensive formatting preservation"""
    
    def __init__(
        self,
        api_key: str,
        model: str = "gpt-4o-mini", #never change it
        translate_notes: bool = False,
        skip_hidden_slides: bool = False
    ):
        """
        Initialize the translator with OpenAI credentials and options.
        
        Args:
            api_key: OpenAI API key
            model: OpenAI model to use
            translate_notes: Whether to translate slide notes
            skip_hidden_slides: Whether to skip hidden slides
        """
        self.model = model
        self.client = OpenAI(api_key=api_key)
        self.translate_notes = translate_notes
        self.skip_hidden_slides = skip_hidden_slides
        logger.info(f"Initialized translator with model: {model}")
        logger.info(f"Translate notes: {translate_notes}")
        logger.info(f"Skip hidden slides: {skip_hidden_slides}")
        
        # Initialize components
        self.formatting_manager = FormattingManager()
        self.enhanced_shape_processor = EnhancedShapeProcessor(self.formatting_manager)
        self.validator = FormatValidator(self.formatting_manager)
        
    def translate_text(self, text: str) -> str:
        """Translate text using OpenAI's API"""
        if not text or not text.strip():
            return text
            
        try:
            logger.info("-" * 80)
            logger.info(f"Original text: {text}")
            
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": INSTRUCTIONS},
                    {"role": "user", "content": f"Translate this Italian text to English, preserving all line breaks and spacing:\n\n{text}"}
                ],
                temperature=0.3,  # Lower temperature for more consistent translations
                max_tokens=1000,
                top_p=1,
                frequency_penalty=0,
                presence_penalty=0
            )
            
            translated = response.choices[0].message.content
            logger.info(f"Translated text: {translated}")
            logger.info("-" * 80)
            
            # Verify translation happened
            if translated.strip() == text.strip():
                logger.warning("WARNING: Text remained unchanged after translation attempt!")
            
            return translated
            
        except Exception as e:
            logger.error(f"Translation error: {str(e)}")
            return text
            
    def translate_presentation(
        self,
        input_path: str,
        output_path: str,
        translation_func: Optional[Callable[[str], str]] = None
    ) -> None:
        """
        Translate PowerPoint presentation while preserving formatting.
        
        Args:
            input_path: Path to the input PowerPoint file
            output_path: Path where the translated file will be saved
            translation_func: Optional custom translation function
        """
        # Use provided translation function or default to OpenAI
        translate = translation_func or self.translate_text
        
        try:
            # Load presentations
            logger.info(f"Loading presentation: {input_path}")
            original_prs = Presentation(input_path)
            translated_prs = Presentation(input_path)  # Create a copy
            
            logger.info("\nStarting translation process...")
            translation_count = 0
            unchanged_count = 0
            
            # Process each slide
            for slide_idx, (orig_slide, trans_slide) in enumerate(
                zip(original_prs.slides, translated_prs.slides), 1
            ):
                # Skip hidden slides if configured
                if self.skip_hidden_slides:
                    # Check if slide is hidden through its properties
                    if hasattr(trans_slide, '_element') and hasattr(trans_slide._element, 'show'):
                        if not trans_slide._element.show:
                            logger.info(f"\nSkipping hidden slide {slide_idx}")
                            continue
                
                logger.info(f"\nProcessing slide {slide_idx}")
                
                # Process shapes
                for shape_idx, shape in enumerate(trans_slide.shapes, 1):
                    logger.info(f"\nProcessing shape {shape_idx} on slide {slide_idx}")
                    
                    # Extract text before translation
                    if hasattr(shape, 'text_frame'):
                        original_text = shape.text_frame.text.strip()
                        if original_text:
                            logger.info(f"Found text to translate: {original_text}")
                    
                    # Process shape with enhanced formatting preservation
                    self.enhanced_shape_processor.process_shape(shape, translate)
                    
                    # Verify translation after processing
                    if hasattr(shape, 'text_frame'):
                        translated_text = shape.text_frame.text.strip()
                        if translated_text:
                            if translated_text == original_text:
                                unchanged_count += 1
                                logger.warning(
                                    f"WARNING: Text remained unchanged: {original_text}"
                                )
                            else:
                                translation_count += 1
                                logger.info(
                                    f"Successfully translated: {original_text} -> "
                                    f"{translated_text}"
                                )
                    
                # Process notes if configured
                if self.translate_notes and trans_slide.has_notes_slide:
                    logger.info(f"\nProcessing notes on slide {slide_idx}")
                    notes_content = TextProcessor.extract_notes_content(
                        trans_slide.notes_slide
                    )
                    if notes_content:
                        text, context_id = notes_content
                        logger.info(f"Found notes to translate: {text}")
                        translated_notes = translate(text)
                        
                        if text.strip() == translated_notes.strip():
                            unchanged_count += 1
                            logger.warning(f"WARNING: Notes remained unchanged: {text}")
                        else:
                            translation_count += 1
                            logger.info(
                                f"Successfully translated notes: {text} -> "
                                f"{translated_notes}"
                            )
                        
                        # Update notes with translation
                        notes_frame = trans_slide.notes_slide.notes_text_frame
                        notes_frame.clear()
                        p = notes_frame.add_paragraph()
                        p.text = translated_notes
            
            logger.info("\nTranslation Statistics:")
            logger.info(f"- Texts translated: {translation_count}")
            logger.info(f"- Texts unchanged: {unchanged_count}")
                        
            # Save translated presentation
            logger.info(f"\nSaving translated presentation to: {output_path}")
            translated_prs.save(output_path)
            logger.info("Translation completed successfully")
            
        except Exception as e:
            logger.error(f"Translation failed: {str(e)}")
            raise

def translate_ppt(
    input_path: str,
    output_path: str,
    api_key: str,
    model: str = "gpt-4o-mini",
    translate_notes: bool = True,
    skip_hidden_slides: bool = True
) -> None:
    """
    Convenience function to translate a PowerPoint file.
    
    Args:
        input_path: Path to the input PowerPoint file
        output_path: Path where the translated file will be saved
        api_key: OpenAI API key
        model: OpenAI model to use (default: gpt-4o-mini)
        translate_notes: Whether to translate slide notes (default: True)
        skip_hidden_slides: Whether to skip hidden slides (default: True)
    """
    translator = PowerPointTranslator(
        api_key,
        model,
        translate_notes=translate_notes,
        skip_hidden_slides=skip_hidden_slides
    )
    translator.translate_presentation(input_path, output_path)




if __name__ == "__main__":
    # Example usage - set your API key in environment variable OPENAI_API_KEY
    import os
    API_KEY = os.getenv('OPENAI_API_KEY')
    
    if not API_KEY:
        print("Please set OPENAI_API_KEY environment variable")
        exit(1)
    
    # Example usage with test file
    translate_ppt(
        input_path='input.pptx',
        output_path='output_translated.pptx',
        api_key=API_KEY,
        translate_notes=False,  # Set to False to skip notes translation
        skip_hidden_slides=False  # Set to False to process hidden slides
    )
