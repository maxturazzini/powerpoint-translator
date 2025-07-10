#!/usr/bin/env python3
import sys
sys.path.append('.')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging

# Disable logging for cleaner output
logging.disable(logging.CRITICAL)

def analyze_presentation(pptx_path):
    """Analyze PowerPoint presentation formatting patterns"""
    try:
        prs = Presentation(pptx_path)
        print(f'✅ Successfully loaded presentation: {pptx_path}')
        print(f'📊 Number of slides: {len(prs.slides)}')
        print('=' * 80)
        
        mixed_format_count = 0
        total_shapes = 0
        shapes_with_text = 0
        complex_formatting_examples = []
        
        # Analyze each slide
        for slide_idx, slide in enumerate(prs.slides, 1):
            print(f'\n🔍 SLIDE {slide_idx} ANALYSIS:')
            print(f'   Number of shapes: {len(slide.shapes)}')
            
            # Analyze each shape
            for shape_idx, shape in enumerate(slide.shapes, 1):
                total_shapes += 1
                print(f'\n   📦 Shape {shape_idx}:')
                print(f'      Type: {shape.shape_type}')
                print(f'      Shape ID: {shape.shape_id}')
                
                # Analyze text content if available
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    shapes_with_text += 1
                    print(f'      Has text frame: Yes')
                    print(f'      Number of paragraphs: {len(shape.text_frame.paragraphs)}')
                    
                    # Analyze each paragraph
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs, 1):
                        print(f'         📝 Paragraph {para_idx}:')
                        print(f'            Number of runs: {len(paragraph.runs)}')
                        print(f'            Alignment: {paragraph.alignment}')
                        print(f'            Level: {paragraph.level}')
                        
                        # Check if paragraph has mixed formatting
                        if len(paragraph.runs) > 1:
                            mixed_format_count += 1
                            print(f'            ⚠️  MIXED FORMATTING DETECTED!')
                            
                            # Create detailed analysis of mixed formatting
                            mixed_details = {
                                'slide': slide_idx,
                                'shape': shape_idx,
                                'paragraph': para_idx,
                                'runs': []
                            }
                            
                            # Analyze each run for mixed formatting
                            for run_idx, run in enumerate(paragraph.runs, 1):
                                font = run.font
                                text_preview = run.text[:30] + ('...' if len(run.text) > 30 else '')
                                print(f'               🔤 Run {run_idx}: "{text_preview}"')
                                
                                # Collect detailed formatting info
                                run_format = {
                                    'text': run.text,
                                    'font_name': font.name,
                                    'font_size': font.size,
                                    'bold': font.bold,
                                    'italic': font.italic,
                                    'underline': font.underline,
                                    'color_type': None,
                                    'language': None
                                }
                                
                                print(f'                  Font: {font.name}')
                                print(f'                  Size: {font.size}')
                                print(f'                  Bold: {font.bold}')
                                print(f'                  Italic: {font.italic}')
                                print(f'                  Underline: {font.underline}')
                                
                                # Check color information
                                if font.color:
                                    try:
                                        color_type = font.color.type
                                        run_format['color_type'] = color_type
                                        print(f'                  Color type: {color_type}')
                                        if hasattr(font.color, 'rgb') and font.color.rgb:
                                            print(f'                  RGB: {font.color.rgb}')
                                        if hasattr(font.color, 'theme_color'):
                                            print(f'                  Theme color: {font.color.theme_color}')
                                    except Exception as e:
                                        print(f'                  Color error: {e}')
                                
                                # Check language
                                try:
                                    lang = font.language_id
                                    run_format['language'] = lang
                                    print(f'                  Language: {lang}')
                                except Exception as e:
                                    print(f'                  Language error: {e}')
                                
                                mixed_details['runs'].append(run_format)
                            
                            complex_formatting_examples.append(mixed_details)
                        
                        # Even single runs can have interesting formatting
                        elif len(paragraph.runs) == 1:
                            run = paragraph.runs[0]
                            font = run.font
                            text_preview = run.text[:30] + ('...' if len(run.text) > 30 else '')
                            print(f'               🔤 Single Run: "{text_preview}"')
                            print(f'                  Font: {font.name}')
                            print(f'                  Size: {font.size}')
                            print(f'                  Bold: {font.bold}')
                            print(f'                  Italic: {font.italic}')
                            print(f'                  Underline: {font.underline}')
                
                # Check for tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    print(f'      📊 TABLE DETECTED')
                    try:
                        table = shape.table
                        print(f'         Rows: {len(table.rows)}')
                        print(f'         Columns: {len(table.columns)}')
                        
                        # Sample some cells
                        for row_idx in range(min(2, len(table.rows))):
                            row = table.rows[row_idx]
                            for col_idx in range(min(2, len(row.cells))):
                                cell = row.cells[col_idx]
                                if cell.text_frame:
                                    cell_text = cell.text_frame.text[:20] + ('...' if len(cell.text_frame.text) > 20 else '')
                                    print(f'         Cell [{row_idx + 1},{col_idx + 1}]: "{cell_text}"')
                                    if len(cell.text_frame.paragraphs) > 0:
                                        para = cell.text_frame.paragraphs[0]
                                        if len(para.runs) > 1:
                                            print(f'            ⚠️  MIXED FORMATTING IN CELL!')
                                            mixed_format_count += 1
                    except Exception as e:
                        print(f'         Table analysis error: {e}')
                
                # Check for groups
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    print(f'      👥 GROUP DETECTED')
                    try:
                        print(f'         Number of shapes in group: {len(shape.shapes)}')
                        
                        # Analyze group shapes
                        for group_shape_idx, group_shape in enumerate(shape.shapes, 1):
                            print(f'            Group Shape {group_shape_idx}: {group_shape.shape_type}')
                            if hasattr(group_shape, 'text_frame') and group_shape.text_frame:
                                if len(group_shape.text_frame.paragraphs) > 0:
                                    para = group_shape.text_frame.paragraphs[0]
                                    if len(para.runs) > 1:
                                        print(f'               ⚠️  MIXED FORMATTING IN GROUP SHAPE!')
                                        mixed_format_count += 1
                    except Exception as e:
                        print(f'         Group analysis error: {e}')
                
                # Check for other shape types
                else:
                    shape_type_name = str(shape.shape_type)
                    print(f'      🔶 Other shape type: {shape_type_name}')
            
            # Check for notes
            if slide.has_notes_slide:
                print(f'\n   📝 NOTES SLIDE DETECTED')
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    notes_preview = notes_text[:50] + ('...' if len(notes_text) > 50 else '')
                    print(f'      Notes text: "{notes_preview}"')
                except Exception as e:
                    print(f'      Notes analysis error: {e}')
            
            print('-' * 60)
        
        # Summary
        print(f'\n📈 FORMATTING ANALYSIS SUMMARY:')
        print(f'   Total slides: {len(prs.slides)}')
        print(f'   Total shapes: {total_shapes}')
        print(f'   Shapes with text: {shapes_with_text}')
        print(f'   Mixed formatting instances: {mixed_format_count}')
        print(f'   Complex formatting examples: {len(complex_formatting_examples)}')
        
        # Display complex formatting examples
        if complex_formatting_examples:
            print(f'\n🎯 COMPLEX FORMATTING EXAMPLES:')
            for idx, example in enumerate(complex_formatting_examples[:3], 1):  # Show first 3
                print(f'\n   Example {idx} (Slide {example["slide"]}, Shape {example["shape"]}, Para {example["paragraph"]}):')
                for run_idx, run_data in enumerate(example['runs'], 1):
                    print(f'      Run {run_idx}: "{run_data["text"][:20]}..."')
                    print(f'         Bold: {run_data["bold"]}, Italic: {run_data["italic"]}, Underline: {run_data["underline"]}')
                    print(f'         Font: {run_data["font_name"]}, Size: {run_data["font_size"]}')
        
        return True
        
    except Exception as e:
        print(f'❌ Error loading presentation: {e}')
        print(f'File path: {pptx_path}')
        return False

if __name__ == "__main__":
    analyze_presentation('sample_pptx/renewable_energy_sample_translation.pptx')
