#!/usr/bin/env python3
import sys
import os
sys.path.append('.')
from translate_powerpoint import PowerPointTranslator
import logging

# Set up logging to see what happens during translation
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def mock_translate(text):
    """Mock translation function for testing - just adds '[ENHANCED]' prefix"""
    if not text or not text.strip():
        return text
    return f"[ENHANCED] {text}"

def test_enhanced_translator():
    """Test the enhanced translator with the sample file"""
    try:
        # Initialize translator with dummy API key
        translator = PowerPointTranslator(
            api_key="dummy_key_for_testing",
            model="gpt-4o-mini",
            translate_notes=True,
            skip_hidden_slides=False
        )
        
        input_path = 'sample_pptx/renewable_energy_sample_translation.pptx'
        output_path = 'sample_pptx/test_enhanced_output.pptx'
        
        print(f"🧪 Testing ENHANCED translator with: {input_path}")
        print(f"📤 Output will be saved to: {output_path}")
        
        # Use mock translation function instead of OpenAI
        translator.translate_presentation(input_path, output_path, mock_translate)
        
        print("✅ Enhanced translation completed successfully!")
        print(f"📁 Check the enhanced output file: {output_path}")
        
        # Now analyze the output to see formatting preservation
        print("\n🔍 Analyzing enhanced output formatting...")
        from sample_pptx.analyze_sample import analyze_presentation
        print("\n" + "="*60)
        print("ANALYZING ENHANCED OUTPUT FILE:")
        print("="*60)
        analyze_presentation(output_path)
        
        # Compare with original
        print("\n" + "="*60)
        print("COMPARISON SUMMARY:")
        print("="*60)
        
        # Load both files for comparison
        from pptx import Presentation
        
        original_prs = Presentation(input_path)
        enhanced_prs = Presentation(output_path)
        
        print("\n🎯 MIXED FORMATTING PRESERVATION TEST:")
        
        # Test Slide 1, Shape 3 (Bold and Italic text)
        orig_shape = original_prs.slides[0].shapes[2]  # Shape 3
        enhanced_shape = enhanced_prs.slides[0].shapes[2]
        
        print(f"\nSlide 1, Shape 3 - Original:")
        for para_idx, para in enumerate(orig_shape.text_frame.paragraphs):
            print(f"  Paragraph {para_idx + 1}: {len(para.runs)} runs")
            for run_idx, run in enumerate(para.runs):
                print(f"    Run {run_idx + 1}: '{run.text}' (Bold: {run.font.bold}, Italic: {run.font.italic})")
        
        print(f"\nSlide 1, Shape 3 - Enhanced:")
        for para_idx, para in enumerate(enhanced_shape.text_frame.paragraphs):
            print(f"  Paragraph {para_idx + 1}: {len(para.runs)} runs")
            for run_idx, run in enumerate(para.runs):
                print(f"    Run {run_idx + 1}: '{run.text}' (Bold: {run.font.bold}, Italic: {run.font.italic})")
        
        # Test Slide 2, Shape 2 (Key Points with bold formatting)
        orig_shape2 = original_prs.slides[1].shapes[1]  # Shape 2
        enhanced_shape2 = enhanced_prs.slides[1].shapes[1]
        
        print(f"\nSlide 2, Shape 2 - Original:")
        for para_idx, para in enumerate(orig_shape2.text_frame.paragraphs):
            if len(para.runs) > 1:  # Only show mixed formatting paragraphs
                print(f"  Paragraph {para_idx + 1}: {len(para.runs)} runs")
                for run_idx, run in enumerate(para.runs):
                    print(f"    Run {run_idx + 1}: '{run.text}' (Bold: {run.font.bold}, Italic: {run.font.italic})")
        
        print(f"\nSlide 2, Shape 2 - Enhanced:")
        for para_idx, para in enumerate(enhanced_shape2.text_frame.paragraphs):
            if len(para.runs) > 1:  # Only show mixed formatting paragraphs
                print(f"  Paragraph {para_idx + 1}: {len(para.runs)} runs")
                for run_idx, run in enumerate(para.runs):
                    print(f"    Run {run_idx + 1}: '{run.text}' (Bold: {run.font.bold}, Italic: {run.font.italic})")
        
        print(f"\n🏆 ENHANCEMENT SUCCESS METRICS:")
        
        # Count mixed formatting preservation
        original_mixed_count = 0
        enhanced_mixed_count = 0
        
        for slide in original_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        if len(para.runs) > 1:
                            original_mixed_count += 1
                            
        for slide in enhanced_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        if len(para.runs) > 1:
                            enhanced_mixed_count += 1
        
        print(f"   Original mixed formatting instances: {original_mixed_count}")
        print(f"   Enhanced mixed formatting instances: {enhanced_mixed_count}")
        
        if enhanced_mixed_count >= original_mixed_count:
            print(f"   ✅ SUCCESS: Mixed formatting preserved!")
        else:
            print(f"   ❌ ISSUE: Mixed formatting partially lost")
            
        preservation_rate = (enhanced_mixed_count / max(original_mixed_count, 1)) * 100
        print(f"   📊 Formatting preservation rate: {preservation_rate:.1f}%")
        
    except Exception as e:
        print(f"❌ Error during enhanced translation test: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    test_enhanced_translator()
