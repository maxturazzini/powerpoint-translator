from typing import List, Dict, Tuple, Optional, Any
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging
import re

logger = logging.getLogger('ppt_translator')

class TextProcessor:
    """Enhanced text processing with format preservation"""
    
    @staticmethod
    def extract_shape_content(shape: BaseShape) -> List[Tuple[str, str]]:
        """
        Extract text content from shapes with their context identifiers.
        Returns list of (text, context_id) tuples.
        """
        content = []
        
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Handle grouped shapes
            for subshape in shape.shapes:
                content.extend(TextProcessor.extract_shape_content(subshape))
                
        elif shape.shape_type == 24:  # MSO_SHAPE_TYPE for SmartArt graphics
            # Handle SmartArt
            if hasattr(shape, "element") and shape.element is not None:
                graphic_data = shape.element.xpath(
                    ".//p:graphic/p:graphicData",
                    namespaces={
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    }
                )
                
                if graphic_data:
                    gd = graphic_data[0]
                    # Extract text from all text elements in SmartArt
                    for idx, t_elem in enumerate(gd.iter(qn('a:t'))):
                        if t_elem.text and t_elem.text.strip():
                            context_id = f"{shape.shape_id}_smartart_{idx}"
                            content.append((t_elem.text, context_id))
                            logger.info(f"Extracted SmartArt text: {t_elem.text}")
                            
        elif isinstance(shape, GraphicFrame) and shape.has_table:
            # Handle tables
            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    if cell.text_frame:
                        context_id = f"{shape.shape_id}_table_{row_idx}_{col_idx}"
                        cell_text = TextProcessor.get_text_frame_content(cell.text_frame)
                        if cell_text.strip():
                            content.append((cell_text, context_id))
                            logger.info(f"Extracted table cell text: {cell_text}")
                            
        elif hasattr(shape, 'text_frame'):
            # Handle regular text frames
            if shape.text_frame:
                text = TextProcessor.get_text_frame_content(shape.text_frame)
                if text.strip():
                    content.append((text, str(shape.shape_id)))
                    logger.info(f"Extracted text frame content: {text}")
                    
        return content

    @staticmethod
    def get_text_frame_content(text_frame) -> str:
        """Extract text content from a text frame preserving structure"""
        paragraphs = []
        for paragraph in text_frame.paragraphs:
            # Get text from each run in the paragraph
            para_text = ""
            for run in paragraph.runs:
                if run.text:
                    # Preserve any existing line breaks within the run
                    para_text += run.text.replace('\v', '\n')
            
            # Only add non-empty paragraphs
            if para_text.strip():
                paragraphs.append(para_text)
                
        # Join paragraphs with newlines, preserving any existing blank lines
        text = "\n".join(paragraphs)
        logger.info(f"Extracted text frame content: {text}")
        return text

    @staticmethod
    def insert_format_markers(text: str) -> Tuple[str, Dict[str, Dict[str, Any]]]:
        """
        Insert format markers into text and track their positions and properties.
        Returns (marked_text, marker_info)
        """
        markers = {}
        lines = text.split('\n')
        marked_lines = []
        
        # Track line break types and spacing
        current_pos = 0
        for idx, line in enumerate(lines):
            # Detect line properties
            leading_spaces = len(line) - len(line.lstrip())
            trailing_spaces = len(line) - len(line.rstrip())
            is_blank = not line.strip()
            
            # Create marker with properties
            marker_id = f"P{idx}"
            marker = f"[{marker_id}]"
            
            # Store marker information
            markers[marker_id] = {
                'position': (current_pos, current_pos + len(line)),
                'properties': {
                    'leading_spaces': leading_spaces,
                    'trailing_spaces': trailing_spaces,
                    'is_blank': is_blank,
                    'original_text': line
                }
            }
            
            # Add marker to line
            marked_line = f"{marker}{line}"
            marked_lines.append(marked_line)
            
            # Update position counter
            current_pos += len(line) + 1  # +1 for newline
            
        marked_text = '\n'.join(marked_lines)
        logger.info(f"Text with format markers: {marked_text}")
        return (marked_text, markers)

    @staticmethod
    def remove_format_markers(text: str, markers: Optional[Dict[str, Dict[str, Any]]] = None) -> str:
        """
        Remove format markers from translated text while preserving formatting.
        If markers dict is provided, uses it to restore original spacing.
        """
        if markers:
            # First remove markers but keep positions
            lines = text.split('\n')
            clean_lines = []
            
            for idx, line in enumerate(lines):
                # Remove marker
                clean_line = re.sub(r'\[P\d+\]', '', line)
                
                # Restore original spacing if available
                marker_id = f"P{idx}"
                if marker_id in markers:
                    props = markers[marker_id]['properties']
                    # Restore leading/trailing spaces
                    if not props['is_blank']:
                        clean_line = ' ' * props['leading_spaces'] + clean_line.strip() + ' ' * props['trailing_spaces']
                
                clean_lines.append(clean_line)
                
            clean_text = '\n'.join(clean_lines)
        else:
            # Simple marker removal if no marker info available
            clean_text = re.sub(r'\[P\d+\]', '', text)
            
        logger.info(f"Text after removing markers: {clean_text}")
        return clean_text

    @staticmethod
    def verify_translation_integrity(
        original_content: List[Tuple[str, str]],
        translated_content: List[Tuple[str, str]]
    ) -> List[str]:
        """
        Verify the integrity of translation.
        Returns list of warning messages.
        """
        warnings = []
        
        if len(original_content) != len(translated_content):
            warnings.append(
                f"Content count mismatch: Original={len(original_content)}, "
                f"Translated={len(translated_content)}"
            )
            
        orig_dict = dict(original_content)
        trans_dict = dict(translated_content)
        
        # Check for missing or extra context IDs
        missing_ids = set(orig_dict.keys()) - set(trans_dict.keys())
        extra_ids = set(trans_dict.keys()) - set(orig_dict.keys())
        
        if missing_ids:
            warnings.append(f"Missing translations for contexts: {missing_ids}")
        if extra_ids:
            warnings.append(f"Extra translations found for contexts: {extra_ids}")
            
        # Check for empty translations
        for context_id, text in trans_dict.items():
            if not text.strip() and orig_dict.get(context_id, '').strip():
                warnings.append(f"Empty translation for context: {context_id}")
                
        return warnings

    @staticmethod
    def extract_notes_content(notes_slide) -> Optional[Tuple[str, str]]:
        """Extract content from notes slide"""
        if notes_slide and notes_slide.notes_text_frame:
            notes_text = TextProcessor.get_text_frame_content(
                notes_slide.notes_text_frame
            )
            if notes_text.strip():
                logger.info(f"Extracted notes text: {notes_text}")
                return (notes_text, "notes")
        return None