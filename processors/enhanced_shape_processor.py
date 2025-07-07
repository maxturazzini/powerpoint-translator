from typing import Optional, List, Dict, Tuple
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.text.text import _Run, _Paragraph
from formatting import FormattingManager
from .text_processor import TextProcessor
import logging
import re

logger = logging.getLogger('ppt_translator')

class EnhancedShapeProcessor:
    """Enhanced shape processor with run-preserving translation"""
    
    def __init__(self, formatting_manager: FormattingManager):
        self.formatting_manager = formatting_manager
        self.text_processor = TextProcessor()
        
    def process_shape(self, shape: BaseShape, translate_func) -> None:
        """
        Process a shape based on its type, preserving formatting at run level.
        """
        try:
            logger.info(f"Processing shape type: {shape.shape_type}")
            
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._process_group_shape(shape, translate_func)
            elif shape.shape_type == 24:  # MSO_SHAPE_TYPE for SmartArt graphics
                self._process_smartart(shape, translate_func)
            elif isinstance(shape, GraphicFrame) and shape.has_table:
                self._process_table(shape, translate_func)
            elif hasattr(shape, 'text_frame'):
                self._process_text_frame_enhanced(shape, translate_func)
        except Exception as e:
            logger.error(f"Error processing shape: {str(e)}")
            raise

    def _process_group_shape(self, group_shape: BaseShape, translate_func) -> None:
        """Process all shapes within a group"""
        logger.info("Processing group shape")
        for shape in group_shape.shapes:
            self.process_shape(shape, translate_func)

    def _process_smartart(self, smartart: BaseShape, translate_func) -> None:
        """Process SmartArt graphics with enhanced format preservation"""
        logger.info("Processing SmartArt")
        if not hasattr(smartart, "element") or smartart.element is None:
            return

        # Extract and store formatting before translation
        content_list = self.text_processor.extract_shape_content(smartart)
        
        for text, context_id in content_list:
            if text.strip():
                # Validate text before translation
                if not self._should_translate_text(text):
                    logger.info(f"SmartArt: Skipping '{text}' (validation failed)")
                    continue
                
                logger.info(f"SmartArt text to translate: {text}")
                
                # Translate without format markers for better quality
                translated_text = translate_func(text)
                logger.info(f"SmartArt translated text: {translated_text}")
                
                # Update the text while preserving formatting
                self._update_smartart_text_enhanced(smartart, translated_text, context_id)

    def _process_table(self, table_frame: GraphicFrame, translate_func) -> None:
        """Process table cells with enhanced format preservation"""
        logger.info("Processing table")
        table = table_frame.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text_frame:
                    logger.info(f"Processing table cell [{row_idx + 1},{col_idx + 1}]")
                    self._translate_text_frame_runs(cell.text_frame, translate_func)

    def _process_text_frame_enhanced(self, shape: BaseShape, translate_func) -> None:
        """Process regular text frames with enhanced run-level preservation"""
        if not shape.text_frame:
            return
            
        # Get text content before translation
        text = self.text_processor.get_text_frame_content(shape.text_frame)
        if not text.strip():
            return
            
        logger.info(f"Text frame content to translate: {text}")
        
        # Use enhanced run-level translation
        self._translate_text_frame_runs(shape.text_frame, translate_func)

    def _translate_text_frame_runs(self, text_frame, translate_func) -> None:
        """
        Translate text frame content while preserving run structure.
        This is the core improvement - no more destructive clearing!
        """
        for paragraph in text_frame.paragraphs:
            self._translate_paragraph_runs(paragraph, translate_func)

    def _translate_paragraph_runs(self, paragraph: _Paragraph, translate_func) -> None:
        """
        Translate each run in a paragraph individually, preserving formatting boundaries.
        """
        if not paragraph.runs:
            return
            
        logger.info(f"Translating paragraph with {len(paragraph.runs)} runs")
        
        # Strategy 1: Translate each run individually
        for run_idx, run in enumerate(paragraph.runs):
            if run.text and run.text.strip():
                original_text = run.text
                
                # Validate text before translation
                if not self._should_translate_text(original_text):
                    logger.info(f"  Run {run_idx + 1}: Skipping '{original_text}' (validation failed)")
                    continue
                
                # Preserve whitespace patterns for proper spacing
                leading_spaces = original_text[:len(original_text) - len(original_text.lstrip())]
                trailing_spaces = original_text[len(original_text.rstrip()):]
                core_text = original_text.strip()
                
                # Translate only the core text content
                translated_core = translate_func(core_text)
                
                # Reconstruct with original spacing
                translated_text = leading_spaces + translated_core + trailing_spaces
                
                # Update the run's text while preserving all formatting and spacing
                run.text = translated_text
                
                logger.info(f"  Run {run_idx + 1}: '{original_text}' -> '{translated_text}'")
            
        logger.info(f"Paragraph translation completed with preserved formatting")

    def _intelligent_run_mapping(self, paragraph: _Paragraph, translated_text: str, translate_func) -> None:
        """
        Alternative strategy: Intelligent mapping of translated text back to runs.
        This handles cases where word boundaries might change during translation.
        """
        if not paragraph.runs:
            return
            
        # Collect original run information
        original_runs = []
        for run in paragraph.runs:
            original_runs.append({
                'text': run.text,
                'formatting': self.formatting_manager.collect_run_formatting(run, paragraph)
            })
        
        # For now, use simple run-by-run translation
        # This can be enhanced later for complex text redistribution
        for run_idx, run in enumerate(paragraph.runs):
            if run.text and run.text.strip():
                run.text = translate_func(run.text)

    def _translate_runs_with_context(self, paragraph: _Paragraph, translate_func) -> None:
        """
        Advanced strategy: Translate with context awareness.
        Sends full paragraph to translator but maps result back to runs.
        """
        if not paragraph.runs:
            return
            
        # Build full paragraph text
        full_text = ""
        run_boundaries = []
        current_pos = 0
        
        for run in paragraph.runs:
            run_text = run.text
            run_boundaries.append({
                'start': current_pos,
                'end': current_pos + len(run_text),
                'text': run_text,
                'run': run
            })
            full_text += run_text
            current_pos += len(run_text)
        
        if not full_text.strip():
            return
            
        # Translate full paragraph for better context
        translated_full = translate_func(full_text)
        
        # For now, fall back to individual run translation
        # This mapping logic can be enhanced for better text distribution
        for boundary in run_boundaries:
            if boundary['text'].strip():
                boundary['run'].text = translate_func(boundary['text'])

    def _update_smartart_text_enhanced(
        self,
        smartart: BaseShape,
        new_text: str,
        context_id: str
    ) -> None:
        """Update SmartArt text with enhanced formatting preservation"""
        if not hasattr(smartart, "element") or smartart.element is None:
            return
            
        text_elements = list(
            smartart.element.iter(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}t'
            )
        )
        
        if text_elements:
            # For simple cases, update the first text element
            # This preserves more formatting than word-by-word splitting
            text_elements[0].text = new_text
            
            # Clear other text elements to avoid duplication
            for elem in text_elements[1:]:
                elem.text = ""

    def _store_paragraph_properties(self, paragraph: _Paragraph) -> Dict:
        """Store all properties of a paragraph"""
        return {
            'alignment': paragraph.alignment,
            'level': paragraph.level
        }

    def _apply_paragraph_properties(self, paragraph: _Paragraph, properties: Dict) -> None:
        """Apply paragraph-level formatting properties"""
        try:
            if 'alignment' in properties and properties['alignment'] is not None:
                paragraph.alignment = properties['alignment']
            if 'level' in properties and properties['level'] is not None:
                paragraph.level = properties['level']
        except Exception as e:
            logger.warning(f"Error applying paragraph properties: {str(e)}")

    def _store_text_frame_properties(self, text_frame) -> Dict:
        """Store all properties of a text frame"""
        return {
            'auto_size': text_frame.auto_size,
            'word_wrap': text_frame.word_wrap,
            'margin_left': text_frame.margin_left,
            'margin_right': text_frame.margin_right,
            'margin_top': text_frame.margin_top,
            'margin_bottom': text_frame.margin_bottom,
            'vertical_anchor': text_frame.vertical_anchor,
        }

    def _apply_text_frame_properties(self, text_frame, properties: Dict) -> None:
        """Apply stored properties to a text frame"""
        try:
            if properties.get('auto_size') is not None:
                text_frame.auto_size = properties['auto_size']
            if properties.get('word_wrap') is not None:
                text_frame.word_wrap = properties['word_wrap']
            if properties.get('margin_left') is not None:
                text_frame.margin_left = properties['margin_left']
            if properties.get('margin_right') is not None:
                text_frame.margin_right = properties['margin_right']
            if properties.get('margin_top') is not None:
                text_frame.margin_top = properties['margin_top']
            if properties.get('margin_bottom') is not None:
                text_frame.margin_bottom = properties['margin_bottom']
            if properties.get('vertical_anchor') is not None:
                text_frame.vertical_anchor = properties['vertical_anchor']
        except Exception as e:
            logger.warning(f"Error applying text frame properties: {str(e)}")
    
    def _should_translate_text(self, text: str) -> bool:
        """
        Validate if text should be translated.
        Skips empty strings, pure whitespace, and uppercase acronyms.
        """
        # Remove leading/trailing whitespace for analysis
        clean_text = text.strip()
        
        # Skip empty strings
        if not clean_text:
            return False
        
        # Skip pure whitespace
        if text.isspace():
            return False
        
        # Skip single characters (likely bullet points, numbers, etc.)
        if len(clean_text) == 1:
            return False
        
        # Skip uppercase acronyms/abbreviations (2+ characters, all uppercase, may contain numbers)
        if len(clean_text) >= 2 and clean_text.isupper() and not clean_text.isdigit():
            # Allow translation if it contains lowercase letters or common punctuation
            if any(c.islower() for c in clean_text):
                return True
            # Skip pure uppercase text that looks like an acronym
            return False
        
        # Skip pure numbers
        if clean_text.isdigit():
            return False
        
        # Allow translation for everything else
        return True