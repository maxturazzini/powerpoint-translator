from typing import Optional, List, Dict
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.text.text import _Run, _Paragraph
from formatting import FormattingManager
from .text_processor import TextProcessor
import logging

logger = logging.getLogger('ppt_translator')

class ShapeProcessor:
    """Process different types of PowerPoint shapes with format preservation"""
    
    def __init__(self, formatting_manager: FormattingManager):
        self.formatting_manager = formatting_manager
        self.text_processor = TextProcessor()
        
    def process_shape(self, shape: BaseShape, translate_func) -> None:
        """
        Process a shape based on its type, preserving formatting.
        """
        try:
            logger.info(f"Processing shape type: {shape.shape_type}")
            
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._process_group_shape(shape, translate_func)
            elif shape.shape_type == 24:  # MSO_SHAPE_TYPE for SmartArt graphics
                self._process_smartart(shape, translate_func)
            elif isinstance(shape, GraphicFrame) and shape.has_table:
                self._process_table(shape, translate_func)
            elif hasattr(shape, 'text_frame'):
                self._process_text_frame(shape, translate_func)
        except Exception as e:
            logger.error(f"Error processing shape: {str(e)}")
            raise

    def _process_group_shape(self, group_shape: BaseShape, translate_func) -> None:
        """Process all shapes within a group"""
        logger.info("Processing group shape")
        for shape in group_shape.shapes:
            self.process_shape(shape, translate_func)

    def _process_smartart(self, smartart: BaseShape, translate_func) -> None:
        """Process SmartArt graphics with format preservation"""
        logger.info("Processing SmartArt")
        if not hasattr(smartart, "element") or smartart.element is None:
            return

        # Extract and store formatting before translation
        content_list = self.text_processor.extract_shape_content(smartart)
        
        for text, context_id in content_list:
            if text.strip():
                logger.info(f"SmartArt text to translate: {text}")
                # Store original formatting
                self.formatting_manager.store_paragraph_formatting(
                    smartart.text_frame.paragraphs[0],  # SmartArt text
                    context_id
                )
                
                # Translate with format markers
                marked_text, markers = self.text_processor.insert_format_markers(text)
                translated_text = translate_func(marked_text)
                clean_text = self.text_processor.remove_format_markers(translated_text)
                logger.info(f"SmartArt translated text: {clean_text}")
                
                # Update the text while preserving formatting
                self._update_smartart_text(smartart, clean_text, context_id)

    def _process_table(self, table_frame: GraphicFrame, translate_func) -> None:
        """Process table cells with format preservation"""
        logger.info("Processing table")
        table = table_frame.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text_frame:
                    context_id = f"{table_frame.shape_id}_table_{row_idx}_{col_idx}"
                    
                    # Store original formatting
                    for paragraph in cell.text_frame.paragraphs:
                        self.formatting_manager.store_paragraph_formatting(
                            paragraph,
                            context_id
                        )
                    
                    # Get cell content
                    cell_text = self.text_processor.get_text_frame_content(
                        cell.text_frame
                    )
                    
                    if cell_text.strip():
                        logger.info(f"Table cell text to translate: {cell_text}")
                        # Translate with format preservation
                        marked_text, markers = self.text_processor.insert_format_markers(
                            cell_text
                        )
                        translated_text = translate_func(marked_text)
                        clean_text = self.text_processor.remove_format_markers(
                            translated_text
                        )
                        logger.info(f"Table cell translated text: {clean_text}")
                        
                        # Update cell content
                        self._update_text_frame(
                            cell.text_frame,
                            clean_text,
                            context_id
                        )

    def _store_paragraph_properties(self, paragraph: _Paragraph) -> Dict:
        """Store all properties of a paragraph"""
        return {
            'alignment': paragraph.alignment,
            'level': paragraph.level
        }

    def _apply_paragraph_properties(self, paragraph: _Paragraph, properties: Dict) -> None:
        """Apply paragraph-level formatting properties"""
        try:
            # Only apply properties that are safe to set through the API
            if 'alignment' in properties and properties['alignment'] is not None:
                paragraph.alignment = properties['alignment']
            if 'level' in properties and properties['level'] is not None:
                paragraph.level = properties['level']
        except Exception as e:
            logger.warning(f"Error applying paragraph properties: {str(e)}")

    def _store_text_frame_properties(self, text_frame) -> Dict:
        """Store all properties of a text frame"""
        return {
            'auto_size': text_frame.auto_size,
            'word_wrap': text_frame.word_wrap,
            'margin_left': text_frame.margin_left,
            'margin_right': text_frame.margin_right,
            'margin_top': text_frame.margin_top,
            'margin_bottom': text_frame.margin_bottom,
            'vertical_anchor': text_frame.vertical_anchor,
            'paragraphs': [
                self._store_paragraph_properties(p) for p in text_frame.paragraphs
            ]
        }

    def _apply_text_frame_properties(self, text_frame, properties: Dict) -> None:
        """Apply stored properties to a text frame"""
        try:
            # Only apply properties that are safe to set through the API
            if properties.get('auto_size') is not None:
                text_frame.auto_size = properties['auto_size']
            if properties.get('word_wrap') is not None:
                text_frame.word_wrap = properties['word_wrap']
            if properties.get('margin_left') is not None:
                text_frame.margin_left = properties['margin_left']
            if properties.get('margin_right') is not None:
                text_frame.margin_right = properties['margin_right']
            if properties.get('margin_top') is not None:
                text_frame.margin_top = properties['margin_top']
            if properties.get('margin_bottom') is not None:
                text_frame.margin_bottom = properties['margin_bottom']
            if properties.get('vertical_anchor') is not None:
                text_frame.vertical_anchor = properties['vertical_anchor']
        except Exception as e:
            logger.warning(f"Error applying text frame properties: {str(e)}")

    def _process_text_frame(self, shape: BaseShape, translate_func) -> None:
        """Process regular text frames with format preservation"""
        if not shape.text_frame:
            return
            
        # Get text content before translation
        text = self.text_processor.get_text_frame_content(shape.text_frame)
        if not text.strip():
            return
            
        logger.info(f"Text frame content to translate: {text}")
            
        # Store text frame and paragraph properties
        text_frame_props = self._store_text_frame_properties(shape.text_frame)
        
        # Store original formatting for each paragraph
        for paragraph in shape.text_frame.paragraphs:
            self.formatting_manager.store_paragraph_formatting(
                paragraph,
                str(shape.shape_id)
            )
        
        # Translate with format preservation
        marked_text, markers = self.text_processor.insert_format_markers(text)
        translated_text = translate_func(marked_text)
        clean_text = self.text_processor.remove_format_markers(translated_text)
        
        logger.info(f"Text frame translated content: {clean_text}")
        
        # Update text frame with translation
        self._update_text_frame(
            shape.text_frame,
            clean_text,
            str(shape.shape_id),
            text_frame_props
        )

    def _update_text_frame(
        self,
        text_frame,
        new_text: str,
        context_id: str,
        text_frame_props: Optional[Dict] = None
    ) -> None:
        """Update text frame content while preserving formatting"""
        if not new_text.strip():
            return
            
        # Store original properties if not provided
        if text_frame_props is None:
            text_frame_props = self._store_text_frame_properties(text_frame)
            
        # Clear existing content
        text_frame.clear()
        
        # Apply text frame properties
        self._apply_text_frame_properties(text_frame, text_frame_props)
        
        # Add new paragraphs
        format_index = 0
        paragraphs = new_text.split('\n')
        for idx, para_text in enumerate(paragraphs):
            if para_text.strip():
                p = text_frame.add_paragraph()
                p.text = para_text.strip()
                
                # Apply original paragraph properties
                if idx < len(text_frame_props['paragraphs']):
                    self._apply_paragraph_properties(
                        p,
                        text_frame_props['paragraphs'][idx]
                    )
                
                # Apply stored formatting
                format_index = self.formatting_manager.apply_paragraph_formatting(
                    p,
                    context_id,
                    format_index
                )

    def _update_smartart_text(
        self,
        smartart: BaseShape,
        new_text: str,
        context_id: str
    ) -> None:
        """Update SmartArt text while preserving formatting"""
        if not hasattr(smartart, "element") or smartart.element is None:
            return
            
        text_elements = list(
            smartart.element.iter(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}t'
            )
        )
        
        # Split text and update elements
        words = new_text.split()
        for idx, elem in enumerate(text_elements):
            if idx < len(words):
                elem.text = words[idx]
                
        # Apply stored formatting
        self.formatting_manager.apply_paragraph_formatting(
            smartart.text_frame.paragraphs[0],
            context_id
        )